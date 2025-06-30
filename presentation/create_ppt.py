#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Blacklist Management System Architecture PowerPoint Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os

def add_gradient_background(slide, color1, color2):
    """슬라이드에 그라디언트 배경 추가"""
    left = top = 0
    width = prs.slide_width
    height = prs.slide_height
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = -45.0
    
    # 그라디언트 색상 설정
    gradient_stops = fill.gradient_stops
    gradient_stops[0].color.rgb = RGBColor(*color1)
    gradient_stops[1].color.rgb = RGBColor(*color2)
    
    # 배경으로 보내기
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    
    return shape

def add_title_subtitle(slide, title, subtitle, title_size=44, subtitle_size=24):
    """제목과 부제목 추가"""
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(title_size)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(subtitle_size)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(16)  # 16:9 비율
prs.slide_height = Inches(9)

# 슬라이드 1: 타이틀
slide1_layout = prs.slide_layouts[5]  # 빈 레이아웃
slide1 = prs.slides.add_slide(slide1_layout)
add_gradient_background(slide1, (30, 41, 59), (126, 34, 206))  # slate-900 to purple-900

# 타이틀 추가
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "Blacklist Management System"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(60)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# 부제목 추가
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "통합 위협 정보 관리 플랫폼"
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
subtitle_frame.paragraphs[0].font.size = Pt(32)
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(196, 181, 253)

# 주요 특징 박스들
features = [
    ("Enterprise Security", "🛡️"),
    ("100만+ IP 처리", "💾"),
    ("실시간 업데이트", "⚡")
]

x_start = 3.5
for i, (text, icon) in enumerate(features):
    box = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_start + i * 3), Inches(6),
        Inches(2.5), Inches(1.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(88, 28, 135)  # purple-800
    box.line.color.rgb = RGBColor(168, 85, 247)  # purple-500
    box.line.width = Pt(2)
    
    text_frame = box.text_frame
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.2)
    
    p = text_frame.add_paragraph()
    p.text = f"{icon} {text}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(233, 213, 255)

# 슬라이드 2: 시스템 아키텍처
slide2_layout = prs.slide_layouts[5]
slide2 = prs.slides.add_slide(slide2_layout)
add_gradient_background(slide2, (17, 24, 39), (31, 41, 55))  # gray-900 to gray-800

# 제목
title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "시스템 아키텍처"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Kubernetes 클러스터 박스
k8s_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1), Inches(2),
    Inches(14), Inches(5.5)
)
k8s_box.fill.gradient()
k8s_box.fill.gradient_angle = 0
gradient_stops = k8s_box.fill.gradient_stops
gradient_stops[0].color.rgb = RGBColor(37, 99, 235)  # blue-600
gradient_stops[1].color.rgb = RGBColor(30, 64, 175)  # blue-800

k8s_text = k8s_box.text_frame
k8s_text.margin_top = Inches(0.3)
p = k8s_text.add_paragraph()
p.text = "Kubernetes Cluster"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# 컴포넌트 박스들
components = [
    ("Application Pods", "Multi Replicas\nAuto-scaling\nHealth Checks", (59, 130, 246)),  # blue-400
    ("Data Layer", "Redis Cache\nSQLite DB\nPVC Storage", (74, 222, 128)),  # green-400
    ("Networking", "Ingress/NodePort\nService Mesh\nLoad Balancing", (168, 85, 247))  # purple-500
]

for i, (title, content, color) in enumerate(components):
    comp_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2 + i * 4.5), Inches(3.5),
        Inches(3.5), Inches(3)
    )
    comp_box.fill.solid()
    comp_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    comp_box.fill.transparency = 0.9  # 10% 불투명도
    
    text_frame = comp_box.text_frame
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.3)
    
    # 제목
    p1 = text_frame.add_paragraph()
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(*color)
    
    # 내용
    p2 = text_frame.add_paragraph()
    p2.text = content
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(100, 100, 100)

# 슬라이드 3: 데이터 플로우
slide3_layout = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(slide3_layout)
add_gradient_background(slide3, (67, 56, 202), (126, 34, 206))  # indigo-900 to purple-900

# 제목
title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "데이터 수집 플로우"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# 데이터 소스
sources = [
    ("REGTECH", "금융보안원", (79, 70, 229)),  # indigo-600
    ("SECUDIUM", "위협 정보", (126, 34, 206))   # purple-700
]

for i, (name, desc, color) in enumerate(sources):
    box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2.5 + i * 2),
        Inches(3), Inches(1.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*color)
    
    text_frame = box.text_frame
    p1 = text_frame.add_paragraph()
    p1.text = name
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = text_frame.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(200, 200, 255)

# 화살표 1
arrow1 = slide3.shapes.add_connector(
    1, Inches(4), Inches(3.25), Inches(2), Inches(0)
)
arrow1.line.color.rgb = RGBColor(96, 165, 250)
arrow1.line.width = Pt(3)

# 처리 시스템
system_box = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6), Inches(2.75),
    Inches(4), Inches(2.5)
)
system_box.fill.solid()
system_box.fill.fore_color.rgb = RGBColor(16, 185, 129)  # green-500

text_frame = system_box.text_frame
p1 = text_frame.add_paragraph()
p1.text = "Blacklist System"
p1.alignment = PP_ALIGN.CENTER
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

p2 = text_frame.add_paragraph()
p2.text = "수집 → 처리 → 저장\n중복 제거 & 검증"
p2.alignment = PP_ALIGN.CENTER
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(220, 252, 231)

# 화살표 2
arrow2 = slide3.shapes.add_connector(
    1, Inches(10), Inches(4), Inches(2), Inches(0)
)
arrow2.line.color.rgb = RGBColor(96, 165, 250)
arrow2.line.width = Pt(3)

# FortiGate
fortigate_box = slide3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(12), Inches(2.75),
    Inches(3), Inches(2.5)
)
fortigate_box.fill.solid()
fortigate_box.fill.fore_color.rgb = RGBColor(245, 158, 11)  # amber-500

text_frame = fortigate_box.text_frame
p1 = text_frame.add_paragraph()
p1.text = "FortiGate"
p1.alignment = PP_ALIGN.CENTER
p1.font.size = Pt(20)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

p2 = text_frame.add_paragraph()
p2.text = "External\nConnector"
p2.alignment = PP_ALIGN.CENTER
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(254, 243, 199)

# 슬라이드 4: 기술 스택
slide4_layout = prs.slide_layouts[5]
slide4 = prs.slides.add_slide(slide4_layout)
add_gradient_background(slide4, (17, 24, 39), (30, 58, 138))  # gray-900 to blue-900

# 제목
title_box = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "기술 스택"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# 기술 스택 박스들
tech_stack = [
    ("Python 3.9", "Flask, Gunicorn\nPandas, Requests", "🐍", (55, 118, 171), (255, 212, 59)),
    ("Kubernetes", "k3s/k8s 1.24+\nAuto-scaling", "☸️", (37, 99, 235), (30, 64, 175)),
    ("Docker", "Multi-stage Build\nAlpine Linux", "🐳", (6, 182, 212), (8, 145, 178)),
    ("Redis", "In-memory Cache\nTTL Management", "📊", (220, 38, 38), (185, 28, 28))
]

for i, (title, desc, icon, color1, color2) in enumerate(tech_stack):
    box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5 + (i % 2) * 7), Inches(2.5 + (i // 2) * 3),
        Inches(6), Inches(2.5)
    )
    
    # 그라디언트 채우기
    box.fill.gradient()
    box.fill.gradient_angle = -45
    gradient_stops = box.fill.gradient_stops
    gradient_stops[0].color.rgb = RGBColor(*color1)
    gradient_stops[1].color.rgb = RGBColor(*color2)
    
    text_frame = box.text_frame
    text_frame.margin_top = Inches(0.3)
    
    # 아이콘과 제목
    p1 = text_frame.add_paragraph()
    p1.text = f"{icon} {title}"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    # 설명
    p2 = text_frame.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(220, 220, 220)

# 슬라이드 5: CI/CD 파이프라인
slide5_layout = prs.slide_layouts[5]
slide5 = prs.slides.add_slide(slide5_layout)
add_gradient_background(slide5, (20, 83, 45), (13, 148, 136))  # green-900 to teal-700

# 제목
title_box = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "CI/CD 파이프라인"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# 파이프라인 단계들
pipeline_steps = [
    ("Git Push", "main branch", "💻", (31, 41, 55)),
    ("GitHub Actions", "Build & Test", "🔧", (55, 65, 81)),
    ("Registry", "registry.jclee.me", "📦", (126, 34, 206)),
    ("Watchtower", "Auto Deploy", "🚀", (34, 197, 94))
]

for i, (title, desc, icon, color) in enumerate(pipeline_steps):
    # 원형 박스
    circle = slide5.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.5 + i * 3.5), Inches(3),
        Inches(2), Inches(2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*color)
    
    # 텍스트
    text_box = slide5.shapes.add_textbox(
        Inches(1.5 + i * 3.5), Inches(5.5),
        Inches(2), Inches(1.5)
    )
    text_frame = text_box.text_frame
    
    p1 = text_frame.add_paragraph()
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    
    p2 = text_frame.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(134, 239, 172)
    
    # 화살표 (마지막 제외)
    if i < len(pipeline_steps) - 1:
        arrow = slide5.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.5 + i * 3.5), Inches(3.75),
            Inches(1), Inches(0.5)
        )
        arrow.fill.gradient()
        gradient_stops = arrow.fill.gradient_stops
        gradient_stops[0].color.rgb = RGBColor(52, 211, 153)
        gradient_stops[1].color.rgb = RGBColor(96, 165, 250)

# 슬라이드 6: 성능 메트릭
slide6_layout = prs.slide_layouts[5]
slide6 = prs.slides.add_slide(slide6_layout)
add_gradient_background(slide6, (88, 28, 135), (219, 39, 119))  # purple-900 to pink-600

# 제목
title_box = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "성능 메트릭"
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# 메트릭 박스들
metrics = [
    ("처리 용량", [("최대 IP 처리", "100만+")], (236, 72, 153)),
    ("응답 시간", [("캐시 히트", "< 10ms"), ("캐시 미스", "< 100ms"), ("대량 조회", "< 500ms")], (250, 204, 21)),
    ("리소스 사용", [("기본 메모리", "~300MB"), ("10만 IP", "~800MB"), ("100만 IP", "~2GB")], (96, 165, 250)),
    ("가용성", [("SLA", "99.9%"), ("자동 복구", "활성화"), ("무중단 배포", "지원")], (74, 222, 128))
]

for i, (title, items, color) in enumerate(metrics):
    box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1 + (i % 2) * 7.5), Inches(2 + (i // 2) * 3.2),
        Inches(6.5), Inches(2.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.fill.transparency = 0.9
    
    text_frame = box.text_frame
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    
    # 제목
    p_title = text_frame.add_paragraph()
    p_title.text = f"📊 {title}"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(*color)
    
    # 항목들
    for label, value in items:
        p_item = text_frame.add_paragraph()
        p_item.text = f"{label}: {value}"
        p_item.font.size = Pt(16)
        p_item.font.color.rgb = RGBColor(75, 85, 99)
        p_item.level = 1

# 파일 저장
output_path = "/home/jclee/app/blacklist/presentation/Blacklist_Architecture.pptx"
prs.save(output_path)
print(f"✅ PowerPoint 파일이 생성되었습니다: {output_path}")

# 추가 정보 출력
print("\n📊 생성된 슬라이드:")
print("1. 타이틀 슬라이드 - 시스템 소개")
print("2. 시스템 아키텍처 - Kubernetes 구조")
print("3. 데이터 플로우 - 수집 및 처리 과정")
print("4. 기술 스택 - 사용 기술 소개")
print("5. CI/CD 파이프라인 - 자동 배포 프로세스")
print("6. 성능 메트릭 - 처리 용량 및 성능 지표")
print("\n💡 파일을 PowerPoint에서 열어 추가 편집이 가능합니다!")